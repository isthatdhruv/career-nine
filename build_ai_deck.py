"""
Builds 'Where We Use AI.pptx' — a presentation for the manager
covering 8 AI applications in the Career-9 platform.

Palette (matches the admin dashboard):
  Base slate:  #0f172a / #1a2238 / #1e293b
  Accent rose: #f43f5e / #fda4af
  Warm amber: #f59e0b
  Slate-500:   #64748b
  Off-white:   #fafafa
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------- Palette ----------
SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
SLATE_800 = RGBColor(0x1E, 0x29, 0x3B)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_400 = RGBColor(0x94, 0xA3, 0xB8)
SLATE_300 = RGBColor(0xCB, 0xD5, 0xE1)
SLATE_100 = RGBColor(0xF1, 0xF5, 0xF9)
OFF_WHITE = RGBColor(0xFA, 0xFA, 0xFA)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

ROSE_500 = RGBColor(0xF4, 0x3F, 0x5E)
ROSE_400 = RGBColor(0xFB, 0x71, 0x85)
ROSE_300 = RGBColor(0xFD, 0xA4, 0xAF)
ROSE_100 = RGBColor(0xFF, 0xE4, 0xE6)

AMBER_500 = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_300 = RGBColor(0xFC, 0xD3, 0x4D)


# ---------- Setup ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    s.shadow.inherit = False
    return s


def add_round_rect(slide, x, y, w, h, fill, radius=0.06, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = radius
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=SLATE_900,
             font="Inter", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    if isinstance(text, list):
        lines = text
    else:
        lines = [text]

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rich_text(slide, x, y, w, h, runs, font="Inter",
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs: list of (text, size, bold, color) per line, OR list of lists for multi-paragraph."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    # normalize: list of paragraphs, each is a list of runs
    if runs and not isinstance(runs[0], list):
        paragraphs = [runs]
    else:
        paragraphs = runs

    for pi, para_runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        for (text, size, bold, color) in para_runs:
            r = p.add_run()
            r.text = text
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
    return tb


def slate_background(slide):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, SLATE_900)
    # Subtle dark band for depth
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), ROSE_500)


def light_background(slide):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, OFF_WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), ROSE_500)


def add_corner_glow(slide, color=ROSE_500, x_frac=0.78, y_frac=-0.15, size=Inches(5)):
    """Decorative oval glow in the corner of slate slides."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        int(SLIDE_W * x_frac), int(SLIDE_H * y_frac),
        size, size,
    )
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    # set transparency to ~78% via XML hack
    sp = s.fill.fore_color._xFill
    alpha_val = "22000"  # ~13% opacity
    # Create alpha element
    fld = sp.find(qn('a:srgbClr'))
    if fld is not None:
        alpha = etree.SubElement(fld, qn('a:alpha'))
        alpha.set('val', alpha_val)
    return s


def add_pill(slide, x, y, text, color=ROSE_500, bg=None, size=11):
    """Small pill-shaped label."""
    if bg is None:
        bg = RGBColor(0x1E, 0x29, 0x3B)  # slate-800
    w = Inches(0.06 * len(text) + 0.6)
    h = Inches(0.32)
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.5
    s.fill.solid()
    s.fill.fore_color.rgb = bg
    s.line.color.rgb = color
    s.line.width = Pt(0.5)
    s.shadow.inherit = False
    add_text(slide, x, y, w, h, text, size=size, bold=True,
             color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return w


def add_section_label(slide, x, y, num_text, label_text, color_num=ROSE_400,
                      color_label=SLATE_400):
    """'01 / SECTION' style label."""
    tb = slide.shapes.add_textbox(x, y, Inches(4), Inches(0.35))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = num_text
    r1.font.name = "Inter"
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = color_num
    r2 = p.add_run()
    r2.text = "  /  " + label_text
    r2.font.name = "Inter"
    r2.font.size = Pt(11)
    r2.font.bold = True
    r2.font.color.rgb = color_label
    return tb


def add_metric_card(slide, x, y, w, h, value, label, accent=ROSE_500,
                    bg=SLATE_800, value_color=WHITE, label_color=SLATE_300):
    """Card with big number + label."""
    add_round_rect(slide, x, y, w, h, bg, radius=0.08)
    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 x + Inches(0.2), y + Inches(0.18),
                                 Inches(0.3), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False

    add_text(slide, x + Inches(0.25), y + Inches(0.36),
             w - Inches(0.5), Inches(0.7),
             value, size=28, bold=True, color=value_color)
    add_text(slide, x + Inches(0.25), y + h - Inches(0.55),
             w - Inches(0.5), Inches(0.4),
             label, size=10, bold=False, color=label_color)


def footer(slide, page_num, total, on_dark=True):
    line_color = SLATE_700 if on_dark else SLATE_300
    text_color = SLATE_400 if on_dark else SLATE_500
    add_rect(slide, Inches(0.5), SLIDE_H - Inches(0.5),
             SLIDE_W - Inches(1), Emu(5000), line_color)
    add_text(slide, Inches(0.5), SLIDE_H - Inches(0.45),
             Inches(6), Inches(0.3),
             "Career-9  ·  Where We Use AI",
             size=9, color=text_color)
    add_text(slide, SLIDE_W - Inches(2.5), SLIDE_H - Inches(0.45),
             Inches(2), Inches(0.3),
             f"{page_num} / {total}",
             size=9, color=text_color, align=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 1 — Title
# ============================================================
slide1 = prs.slides.add_slide(BLANK)
slate_background(slide1)
add_corner_glow(slide1, ROSE_500, x_frac=0.62, y_frac=-0.4, size=Inches(8))
add_corner_glow(slide1, ROSE_500, x_frac=-0.15, y_frac=0.55, size=Inches(7))

# Eyebrow
add_text(slide1, Inches(0.9), Inches(1.6), Inches(8), Inches(0.4),
         "INTERNAL  ·  PRODUCT  ·  AI APPLICATIONS",
         size=12, bold=True, color=ROSE_300)

# Title
add_text(slide1, Inches(0.9), Inches(2.05), Inches(11), Inches(2.0),
         "Where We Use AI",
         size=72, bold=True, color=WHITE)

# Subtitle
add_text(slide1, Inches(0.9), Inches(3.5), Inches(11), Inches(0.6),
         "Eight ways AI powers Career-9's career-guidance platform",
         size=22, color=SLATE_300)

# Divider line
sep = slide1.shapes.add_connector(1, Inches(0.9), Inches(4.4),
                                  Inches(2.2), Inches(4.4))
sep.line.color.rgb = ROSE_500
sep.line.width = Pt(2.5)

# Caption row
add_text(slide1, Inches(0.9), Inches(4.6), Inches(12), Inches(0.4),
         "Indian context  ·  Behavioural data  ·  Longitudinal cohorts  ·  Career outcomes",
         size=13, bold=True, color=SLATE_400)

# Bottom-right badge
add_round_rect(slide1, SLIDE_W - Inches(3.2), SLIDE_H - Inches(1.4),
               Inches(2.6), Inches(0.7), SLATE_800, radius=0.4)
add_text(slide1, SLIDE_W - Inches(3.2), SLIDE_H - Inches(1.4),
         Inches(2.6), Inches(0.7),
         "Career-9  ·  AI-First",
         size=12, bold=True, color=ROSE_300,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# SLIDE 2 — Agenda / Eight applications grid
# ============================================================
slide2 = prs.slides.add_slide(BLANK)
slate_background(slide2)
add_corner_glow(slide2, ROSE_500, x_frac=0.85, y_frac=-0.25, size=Inches(5))

add_section_label(slide2, Inches(0.9), Inches(0.6), "00", "AGENDA")
add_text(slide2, Inches(0.9), Inches(1.0), Inches(11), Inches(0.9),
         "Eight applications, one platform",
         size=36, bold=True, color=WHITE)
add_text(slide2, Inches(0.9), Inches(1.85), Inches(11), Inches(0.4),
         "Each one solves a problem generic career tools cannot.",
         size=14, color=SLATE_400)

# 4x2 grid of cards
items = [
    ("01", "Career pathway prediction",   "Fine-tuned on 250+ careers"),
    ("02", "Personalised insights",        "Potential × values → fit"),
    ("03", "Behavioural proctoring AI",    "Attention, focus, gaze"),
    ("04", "Indian-context model",         "Multilingual, validity-preserved"),
    ("05", "Quality & bias detection",     "NLP on open responses"),
    ("06", "Success probability",          "Longitudinal cohort data"),
    ("07", "AI counsellors",               "On-demand guidance at scale"),
    ("08", "Gamified mobile app",          "Duolingo-style learning"),
]

card_w = Inches(2.85)
card_h = Inches(1.7)
gap = Inches(0.18)
start_x = Inches(0.9)
start_y = Inches(2.65)

for idx, (num, title, sub) in enumerate(items):
    row, col = idx // 4, idx % 4
    x = start_x + col * (card_w + gap)
    y = start_y + row * (card_h + gap)
    add_round_rect(slide2, x, y, card_w, card_h, SLATE_800, radius=0.07)
    # Top accent
    bar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x + Inches(0.25), y + Inches(0.2),
                                  Inches(0.4), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ROSE_500
    bar.line.fill.background()
    bar.shadow.inherit = False

    add_text(slide2, x + Inches(0.25), y + Inches(0.32),
             card_w - Inches(0.5), Inches(0.32),
             num, size=12, bold=True, color=ROSE_400)
    add_text(slide2, x + Inches(0.25), y + Inches(0.62),
             card_w - Inches(0.5), Inches(0.5),
             title, size=14, bold=True, color=WHITE)
    add_text(slide2, x + Inches(0.25), y + Inches(1.1),
             card_w - Inches(0.5), Inches(0.5),
             sub, size=10.5, color=SLATE_400)

footer(slide2, 2, 11)


# ============================================================
# Generic content slide builder for slides 3-10
# ============================================================
def content_slide(num, total_pages, section, headline, blurb, bullets,
                  metric_left=None, metric_right=None):
    """
    section: '01' style number
    headline: big H1
    blurb: 1-2 sentence caption under headline
    bullets: list of (icon_text, bold_label, body_text)
    metric_left/right: optional (value, label) tuples for KPI cards on the right
    """
    s = prs.slides.add_slide(BLANK)
    light_background(s)

    # Slim left rail with section
    add_rect(s, 0, 0, Inches(0.06), SLIDE_H, ROSE_500)

    # Section label
    add_section_label(s, Inches(0.9), Inches(0.55),
                      f"{section}  /  WHERE WE USE AI",
                      "",
                      color_num=ROSE_500, color_label=SLATE_500)

    # Headline
    add_text(s, Inches(0.9), Inches(1.0), Inches(8.5), Inches(1.4),
             headline, size=34, bold=True, color=SLATE_900)

    # Blurb
    add_text(s, Inches(0.9), Inches(2.45), Inches(8.5), Inches(1.0),
             blurb, size=14, color=SLATE_500)

    # Divider
    sep = s.shapes.add_connector(1, Inches(0.9), Inches(3.8),
                                 Inches(1.8), Inches(3.8))
    sep.line.color.rgb = ROSE_500
    sep.line.width = Pt(2.5)

    # Bullets
    by = Inches(4.0)
    for (icon, label, body) in bullets:
        # Icon dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(0.9), by + Inches(0.06),
                                 Inches(0.32), Inches(0.32))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ROSE_100
        dot.line.fill.background()
        dot.shadow.inherit = False
        add_text(s, Inches(0.9), by + Inches(0.06),
                 Inches(0.32), Inches(0.32),
                 icon, size=10, bold=True, color=ROSE_500,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Bold label + body
        add_rich_text(s, Inches(1.4), by, Inches(8.0), Inches(0.9),
                      [
                          [(label, 13, True, SLATE_900),
                           ("  " + body, 13, False, SLATE_700)]
                      ])
        by += Inches(0.95)

    # Right-side metric panel
    if metric_left or metric_right:
        panel_x = Inches(9.7)
        panel_y = Inches(1.0)
        panel_w = Inches(3.1)
        panel_h = Inches(5.7)
        add_round_rect(s, panel_x, panel_y, panel_w, panel_h, SLATE_900,
                       radius=0.05)
        add_corner_glow(s, ROSE_500, x_frac=0.86, y_frac=0.05, size=Inches(2.5))

        add_text(s, panel_x + Inches(0.3), panel_y + Inches(0.3),
                 panel_w - Inches(0.6), Inches(0.3),
                 "AT A GLANCE", size=10, bold=True, color=ROSE_300)

        if metric_left:
            v, l = metric_left
            add_text(s, panel_x + Inches(0.3), panel_y + Inches(0.85),
                     panel_w - Inches(0.6), Inches(1.1),
                     v, size=44, bold=True, color=WHITE)
            add_text(s, panel_x + Inches(0.3), panel_y + Inches(2.0),
                     panel_w - Inches(0.6), Inches(0.6),
                     l, size=12, color=SLATE_300)

        if metric_right:
            v, l = metric_right
            add_text(s, panel_x + Inches(0.3), panel_y + Inches(3.1),
                     panel_w - Inches(0.6), Inches(1.1),
                     v, size=44, bold=True, color=ROSE_300)
            add_text(s, panel_x + Inches(0.3), panel_y + Inches(4.25),
                     panel_w - Inches(0.6), Inches(0.6),
                     l, size=12, color=SLATE_300)

        # Footer pill in panel
        add_text(s, panel_x + Inches(0.3), panel_y + panel_h - Inches(0.7),
                 panel_w - Inches(0.6), Inches(0.4),
                 "AI-driven  ·  India-trained",
                 size=10, bold=True, color=SLATE_400)

    footer(s, num, total_pages, on_dark=False)
    return s


# ============================================================
# SLIDE 3 — Career pathway prediction
# ============================================================
content_slide(
    num=3, total_pages=11, section="01",
    headline="9 careers from a sea of 250+",
    blurb="Our fine-tuned model selects the nine career pathways most suitable "
          "for each student — with a quantified suitability score per pathway.",
    bullets=[
        ("F", "Fine-tuned base model.",
         "Trained on a curated career dataset of 250+ pathways merged with our internal assessment data."),
        ("R", "Ranked recommendations.",
         "Outputs a top-9 short-list — not a generic single career label."),
        ("S", "Suitability score.",
         "Each pathway carries a confidence score so students see fit, not just fit-or-not."),
        ("4", "Drives the 4-pager report.",
         "Powers the latest 4-page student report consumed by parents and counsellors."),
    ],
    metric_left=("250+", "career pathways modelled"),
    metric_right=("9", "personalised recommendations per student"),
)


# ============================================================
# SLIDE 4 — Personalised dashboard insights
# ============================================================
content_slide(
    num=4, total_pages=11, section="02",
    headline="Potential × values → field of work",
    blurb="The student dashboard surfaces insights tied to each suitable pathway — "
          "linking who they are to where they will thrive.",
    bullets=[
        ("P", "Potential signals.",
         "Cognitive performance, learnability, and behavioural traits feed the model."),
        ("V", "Values signals.",
         "What a student cares about — autonomy, impact, security — becomes a first-class predictor."),
        ("→", "Field-of-work prediction.",
         "Combined potential + values signals predict the field where the student is likely to flourish."),
        ("D", "Live in the dashboard.",
         "Students see why a pathway suits them, not just that it does."),
    ],
    metric_left=("2", "signal families combined"),
    metric_right=("1:1", "insights mapped to each pathway"),
)


# ============================================================
# SLIDE 5 — Behavioural proctoring AI
# ============================================================
content_slide(
    num=5, total_pages=11, section="03",
    headline="What attention reveals about fit",
    blurb="We capture rich behavioural telemetry during assessments — and train a model "
          "to extract focus and attention signals that improve career prediction.",
    bullets=[
        ("M", "Mouse dynamics.",
         "Click cadence, hover paths and movement velocity captured per question."),
        ("E", "Eye-gaze tracking.",
         "Fixation points and gaze sequences logged across the assessment."),
        ("A", "Attention-span signal.",
         "ML extracts focus duration, distraction events and re-engagement patterns."),
        ("→", "Feeds the prediction model.",
         "These behavioural features go into the suitability prediction — not just the answers themselves."),
    ],
    metric_left=("4+", "behavioural channels captured"),
    metric_right=("∞", "data points per assessment"),
)


# ============================================================
# SLIDE 6 — Indian context specialisation
# ============================================================
content_slide(
    num=6, total_pages=11, section="04",
    headline="Built for India, not retrofitted",
    blurb="Trained on Indian students and Indian careers, with psychometric validity "
          "preserved across multiple Indian languages.",
    bullets=[
        ("IN", "Indian dataset.",
         "Models trained on Indian students taking real Indian career pathways — not Western proxies."),
        ("भा", "Multilingual coverage.",
         "Same assessment available across most Indian languages used in schools."),
        ("✓", "Construct validity preserved.",
         "The same psychometric construct measured across translations — not free-form rewriting."),
        ("R", "Local relevance.",
         "Career options reflect what Indian students actually choose, not a Silicon Valley list."),
    ],
    metric_left=("8+", "languages with preserved validity"),
    metric_right=("100%", "career list rooted in Indian pathways"),
)


# ============================================================
# SLIDE 7 — Data quality & bias detection
# ============================================================
content_slide(
    num=7, total_pages=11, section="05",
    headline="Catching the noise foreign tools miss",
    blurb="NLP and pattern recognition flag bad data and Indian-specific response biases — "
          "so the prediction stands on clean ground.",
    bullets=[
        ("F", "Response fatigue.",
         "Drop-offs, rushed clicks and end-of-assessment slumps detected automatically."),
        ("R", "Random-clicking detection.",
         "Pattern recognition surfaces students who clicked through without reading."),
        ("S", "Social desirability bias.",
         "NLP on open-ended answers flags responses tuned for what 'sounds right'."),
        ("C", "Indian cultural patterns.",
         "Response habits unique to Indian students that Western tools never learned to handle."),
    ],
    metric_left=("4", "bias families detected"),
    metric_right=("Auto", "psychometric integrity guard"),
)


# ============================================================
# SLIDE 8 — Career success probability
# ============================================================
content_slide(
    num=8, total_pages=11, section="06",
    headline="Will they thrive — not just fit",
    blurb="Beyond match, our model predicts the probability of staying, growing and "
          "switching — trained on longitudinal data nobody else in India has.",
    bullets=[
        ("S", "Satisfaction probability.",
         "How likely a student is to enjoy the pathway over time — not just on day one."),
        ("R", "Retention probability.",
         "Likelihood of staying in the pathway through inflection points."),
        ("X", "Switching probability.",
         "Risk of pivoting to a different pathway — surfaced before it happens."),
        ("L", "Longitudinal cohorts.",
         "Trained on students we've re-assessed over time — proprietary to Career-9."),
    ],
    metric_left=("3", "outcome probabilities per pathway"),
    metric_right=("Yrs", "longitudinal cohort horizon"),
)


# ============================================================
# SLIDE 9 — AI counsellors
# ============================================================
content_slide(
    num=9, total_pages=11, section="07",
    headline="Counselling at the scale of every student",
    blurb="An AI counsellor that knows the student's profile, suitability scores and "
          "behavioural signals — available the moment a question is asked.",
    bullets=[
        ("Q", "Always available.",
         "On-demand guidance — no wait list, no scheduling, no geography."),
        ("P", "Personalised context.",
         "Anchored on the student's actual assessment + dashboard data, not a generic chatbot."),
        ("E", "Escalation built-in.",
         "Hands off to a human counsellor when the question goes beyond the model's confidence."),
        ("S", "Scales with the user base.",
         "One model, every student — without scaling counsellor headcount linearly."),
    ],
    metric_left=("24/7", "available to every student"),
    metric_right=("1:N", "one model, infinite conversations"),
)


# ============================================================
# SLIDE 10 — Gamified mobile app (Duolingo-style)
# ============================================================
content_slide(
    num=10, total_pages=11, section="08",
    headline="Career growth, gamified",
    blurb="A Duolingo-style mobile experience where students learn about — and progress "
          "toward — their career through bite-sized AI-personalised tasks.",
    bullets=[
        ("D", "Daily streaks.",
         "Short tasks that compound into real career-readiness over weeks and months."),
        ("L", "Personalised lessons.",
         "AI tailors content to the student's recommended pathways and weak areas."),
        ("G", "Game mechanics.",
         "Levels, XP, badges and friends — turning career exploration into a habit."),
        ("M", "Measured progress.",
         "Each task gives the model fresh signal — feeding back into suitability and growth predictions."),
    ],
    metric_left=("Daily", "engagement loop"),
    metric_right=("Bi-direct.", "tasks teach the student AND the model"),
)


# ============================================================
# SLIDE 11 — Why this matters / closing
# ============================================================
slide11 = prs.slides.add_slide(BLANK)
slate_background(slide11)
add_corner_glow(slide11, ROSE_500, x_frac=0.78, y_frac=-0.3, size=Inches(7))
add_corner_glow(slide11, ROSE_500, x_frac=-0.2, y_frac=0.5, size=Inches(6))

add_section_label(slide11, Inches(0.9), Inches(0.55), "09", "WHY IT MATTERS")
add_text(slide11, Inches(0.9), Inches(0.95), Inches(11.5), Inches(1.2),
         "We're not adding AI on top.",
         size=46, bold=True, color=WHITE)
add_text(slide11, Inches(0.9), Inches(1.95), Inches(11.5), Inches(1.0),
         "AI is the platform's spine — from prediction to behaviour to outcomes.",
         size=20, color=SLATE_300)

# Three pillar cards
pillars = [
    ("Built for India",
     "Indian students. Indian careers. Indian languages. Validity preserved across all of them."),
    ("Data nobody else has",
     "Longitudinal cohorts and behavioural telemetry give us a moat generic tools cannot replicate."),
    ("From fit to thriving",
     "Other tools tell students what to pick. We tell them how likely they are to grow inside the pick."),
]

card_w = Inches(3.85)
card_h = Inches(2.5)
gap = Inches(0.25)
total_w = 3 * card_w + 2 * gap
start_x = (SLIDE_W - total_w) // 2
y = Inches(3.6)

for i, (title, body) in enumerate(pillars):
    x = start_x + i * (card_w + gap)
    add_round_rect(slide11, x, y, card_w, card_h, SLATE_800, radius=0.06)
    # Top accent
    bar = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   x + Inches(0.3), y + Inches(0.25),
                                   Inches(0.5), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ROSE_500
    bar.line.fill.background()
    bar.shadow.inherit = False

    add_text(slide11, x + Inches(0.3), y + Inches(0.5),
             card_w - Inches(0.6), Inches(0.6),
             title, size=20, bold=True, color=WHITE)
    add_text(slide11, x + Inches(0.3), y + Inches(1.15),
             card_w - Inches(0.6), Inches(1.3),
             body, size=12, color=SLATE_300)

# Closing CTA
add_text(slide11, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.5),
         "→  Where we go next:  ship AI counsellors,  release the gamified app,  expand cohorts.",
         size=14, bold=True, color=ROSE_300)

footer(slide11, 11, 11)


# ---------- Save ----------
out = r"c:\Users\pop\Desktop\career-nine\career-nine\Where-We-Use-AI.pptx"
prs.save(out)
print("WROTE", out)
