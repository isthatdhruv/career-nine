"""
Builds 'Career-9-Comparison.pptx' — a 5-slide comparison deck pitching
Career-9 against six other career guidance platforms.

Source: Career-9 Comparison.pdf

Palette (slate + emerald + cyan; deliberately no purple/pink):
  Slate:   #0f172a / #1e293b / #475569 / #64748b / #cbd5e1 / #f1f5f9
  Emerald: #059669 / #10b981 / #d1fae5  (Career-9 brand highlight)
  Cyan:    #0891b2 / #06b6d4
  Amber:   #f59e0b                       (call-outs)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ---------- Palette ----------
SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
SLATE_800 = RGBColor(0x1E, 0x29, 0x3B)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_600 = RGBColor(0x47, 0x55, 0x69)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_400 = RGBColor(0x94, 0xA3, 0xB8)
SLATE_300 = RGBColor(0xCB, 0xD5, 0xE1)
SLATE_200 = RGBColor(0xE2, 0xE8, 0xF0)
SLATE_100 = RGBColor(0xF1, 0xF5, 0xF9)
SLATE_50  = RGBColor(0xF8, 0xFA, 0xFC)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

EMERALD_700 = RGBColor(0x04, 0x78, 0x57)
EMERALD_600 = RGBColor(0x05, 0x96, 0x69)
EMERALD_500 = RGBColor(0x10, 0xB9, 0x81)
EMERALD_100 = RGBColor(0xD1, 0xFA, 0xE5)
EMERALD_50  = RGBColor(0xEC, 0xFD, 0xF5)

CYAN_700 = RGBColor(0x0E, 0x74, 0x90)
CYAN_600 = RGBColor(0x08, 0x91, 0xB2)
CYAN_500 = RGBColor(0x06, 0xB6, 0xD4)

AMBER_500 = RGBColor(0xF5, 0x9E, 0x0B)


# ---------- Setup ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def add_rect(slide, left, top, width, height, fill, line=None, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    if not shadow:
        shape.shadow.inherit = False
    return shape


def add_round(slide, left, top, width, height, fill, line=None, radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(
    slide, left, top, width, height,
    text, *,
    size=14, bold=False, color=SLATE_800,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    font="Calibri", line_spacing=1.15,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return box


def add_bg(slide, color):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)


def add_page_header(slide, eyebrow, title, *, eyebrow_color=EMERALD_600):
    # Top accent bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.16), EMERALD_500)

    add_text(
        slide,
        Inches(0.6), Inches(0.45), Inches(12), Inches(0.3),
        eyebrow,
        size=11, bold=True, color=eyebrow_color,
        font="Calibri",
    )
    add_text(
        slide,
        Inches(0.6), Inches(0.75), Inches(12), Inches(0.7),
        title,
        size=30, bold=True, color=SLATE_900,
        font="Calibri",
    )
    # Underline accent
    add_rect(slide, Inches(0.6), Inches(1.55), Inches(0.6), Inches(0.05), EMERALD_500)


def add_footer(slide, page_num, total):
    add_text(
        slide,
        Inches(0.6), Inches(7.05), Inches(6), Inches(0.3),
        "Career-9  |  Comparison Deck",
        size=9, color=SLATE_400,
    )
    add_text(
        slide,
        Inches(7), Inches(7.05), Inches(5.733), Inches(0.3),
        f"{page_num} / {total}",
        size=9, color=SLATE_400, align=PP_ALIGN.RIGHT,
    )


# ===================================================================
# Slide 1 — Title
# ===================================================================
def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, SLATE_900)

    # Decorative emerald gradient bars (left edge)
    add_rect(s, 0, 0, Inches(0.35), SLIDE_H, EMERALD_500)
    add_rect(s, Inches(0.42), 0, Inches(0.08), SLIDE_H, EMERALD_700)

    # Subtle decorative circles
    c1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(0.6), Inches(2.2), Inches(2.2))
    c1.fill.solid(); c1.fill.fore_color.rgb = SLATE_800
    c1.line.fill.background(); c1.shadow.inherit = False

    c2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.4), Inches(5.2), Inches(1.4), Inches(1.4))
    c2.fill.solid(); c2.fill.fore_color.rgb = EMERALD_700
    c2.line.fill.background(); c2.shadow.inherit = False

    # Eyebrow
    add_text(
        s, Inches(1.2), Inches(2.0), Inches(11), Inches(0.4),
        "CAREER GUIDANCE PLATFORM COMPARISON",
        size=13, bold=True, color=EMERALD_500,
    )

    # Main title
    add_text(
        s, Inches(1.2), Inches(2.5), Inches(11), Inches(1.4),
        "Career-9 vs the Industry",
        size=58, bold=True, color=WHITE, line_spacing=1.0,
    )

    # Subtitle
    add_text(
        s, Inches(1.2), Inches(3.9), Inches(11), Inches(0.7),
        "Why a holistic, validated, action-oriented approach wins.",
        size=22, color=SLATE_300, line_spacing=1.2,
    )

    # Tagline pills
    pills = [
        ("Holistic", EMERALD_500),
        ("Actionable", CYAN_500),
        ("Validated", AMBER_500),
    ]
    px = Inches(1.2)
    py = Inches(5.2)
    for label, color in pills:
        pill = add_round(s, px, py, Inches(1.7), Inches(0.55), SLATE_800, radius=0.5)
        # Color dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, px + Inches(0.2), py + Inches(0.2), Inches(0.18), Inches(0.18))
        dot.fill.solid(); dot.fill.fore_color.rgb = color
        dot.line.fill.background(); dot.shadow.inherit = False
        add_text(s, px + Inches(0.45), py + Inches(0.07), Inches(1.3), Inches(0.4),
                 label, size=13, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        px += Inches(1.85)

    # Bottom info line
    add_text(
        s, Inches(1.2), Inches(6.6), Inches(11), Inches(0.4),
        "7 platforms compared  •  10 dimensions  •  prepared for partner schools",
        size=12, color=SLATE_400,
    )


# ===================================================================
# Slide 2 — Why Career-9 Works (5 pillars)
# ===================================================================
def slide_why_career9():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, SLATE_50)
    add_page_header(s, "WHY CAREER-9 WORKS", "Five pillars that set us apart")

    pillars = [
        ("01", "Holistic Profiling",
         "Beyond aptitude — emotional, cognitive, and motivational traits combined."),
        ("02", "Actionable Insights",
         "Reports don't just inform. They guide behaviour through tasks and routines."),
        ("03", "Family Partnership",
         "Built-in plans foster collaboration between teachers, parents and students."),
        ("04", "Validated & Trusted",
         "Backed by leading institutions (NIMHANS, UGC, MIT) and peer-reviewed research."),
        ("05", "Behavioural Impact",
         "An LMS that helps students build real habits aligned with career goals."),
    ]

    # 5 cards in a row
    card_w = Inches(2.3)
    card_h = Inches(4.4)
    gap = Inches(0.18)
    total_w = card_w * 5 + gap * 4
    start_x = (SLIDE_W - total_w) / 2
    y = Inches(2.0)

    for i, (num, title, body) in enumerate(pillars):
        x = start_x + (card_w + gap) * i

        # Card background
        add_round(s, x, y, card_w, card_h, WHITE, radius=0.06)

        # Top emerald header strip
        header_h = Inches(0.7)
        add_round(s, x, y, card_w, header_h, EMERALD_600, radius=0.06)
        # Mask the bottom corners of the strip so only top is rounded
        add_rect(s, x, y + Inches(0.35), card_w, Inches(0.35), EMERALD_600)

        # Number badge
        add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(1.0), Inches(0.4),
                 num, size=18, bold=True, color=WHITE)
        # Small tag
        add_text(s, x + Inches(0.2), y + Inches(0.45), Inches(2.0), Inches(0.25),
                 "PILLAR", size=8, bold=True, color=EMERALD_100)

        # Title
        add_text(s, x + Inches(0.25), y + Inches(0.95), card_w - Inches(0.5), Inches(0.85),
                 title, size=17, bold=True, color=SLATE_900, line_spacing=1.1)

        # Divider
        add_rect(s, x + Inches(0.25), y + Inches(1.85), Inches(0.4), Inches(0.04), EMERALD_500)

        # Body
        add_text(s, x + Inches(0.25), y + Inches(2.05), card_w - Inches(0.5), Inches(2.2),
                 body, size=12, color=SLATE_600, line_spacing=1.35)


# ===================================================================
# Comparison-table helpers
# ===================================================================
PLATFORMS = [
    "Career-9",
    "Mindler",
    "iDream\nCareer",
    "Dheya",
    "Mentoria",
    "Pearson\nMCMF",
    "Tamanna\n(Govt.)",
]

# Geometry shared across both comparison slides
TABLE_LEFT   = Inches(0.6)
TABLE_TOP    = Inches(1.85)
FEATURE_COL_W = Inches(1.85)
PLATFORM_COL_W = Inches(1.51)   # 7 cols * 1.51 = 10.57
HEADER_H = Inches(0.55)
ROW_H_DEFAULT = Inches(1.05)


def draw_table_header(slide, top):
    x = TABLE_LEFT
    # Feature header
    add_round(slide, x, top, FEATURE_COL_W, HEADER_H, SLATE_900, radius=0.15)
    add_text(slide, x + Inches(0.15), top, FEATURE_COL_W - Inches(0.3), HEADER_H,
             "FEATURE", size=11, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT)

    x += FEATURE_COL_W + Inches(0.05)
    for i, name in enumerate(PLATFORMS):
        is_career9 = (i == 0)
        bg = EMERALD_600 if is_career9 else SLATE_700
        add_round(slide, x, top, PLATFORM_COL_W, HEADER_H, bg, radius=0.18)
        add_text(slide, x, top, PLATFORM_COL_W, HEADER_H,
                 name, size=10, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, line_spacing=1.0)
        x += PLATFORM_COL_W + Inches(0.05)


def draw_table_row(slide, top, feature_label, values, row_h=ROW_H_DEFAULT, highlight_career9=True):
    x = TABLE_LEFT
    # Feature label cell
    add_round(slide, x, top, FEATURE_COL_W, row_h, SLATE_100, radius=0.1)
    add_text(slide, x + Inches(0.18), top, FEATURE_COL_W - Inches(0.36), row_h,
             feature_label, size=11, bold=True, color=SLATE_900,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)

    x += FEATURE_COL_W + Inches(0.05)
    for i, val in enumerate(values):
        is_career9 = (i == 0 and highlight_career9)
        bg = EMERALD_50 if is_career9 else WHITE
        line = EMERALD_500 if is_career9 else SLATE_200
        cell = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, PLATFORM_COL_W, row_h)
        cell.adjustments[0] = 0.1
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        cell.line.color.rgb = line
        cell.line.width = Pt(1.25 if is_career9 else 0.5)
        cell.shadow.inherit = False

        add_text(
            slide, x + Inches(0.1), top + Inches(0.08),
            PLATFORM_COL_W - Inches(0.2), row_h - Inches(0.16),
            val,
            size=8.5, bold=is_career9,
            color=EMERALD_700 if is_career9 else SLATE_600,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=1.15,
        )
        x += PLATFORM_COL_W + Inches(0.05)


# ===================================================================
# Slide 3 — Methodology & Output
# ===================================================================
def slide_compare_methodology():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, SLATE_50)
    add_page_header(s, "PART 1 OF 2  •  METHODOLOGY & OUTPUT",
                    "How each platform measures and reports")

    draw_table_header(s, TABLE_TOP)

    rows = [
        (
            "Core Framework",
            [
                "6D: Values, Personality, Intelligence, Aspirations, Ability, Knowledge",
                "5D: Orientation, Interests, Personality, Aptitude, EQ",
                "General psychometrics",
                "7D: Direction, Discovery, Development, etc.",
                "4-step: Discovery, Matching, Mentoring, Placement",
                "12 Personality Traits",
                "Aptitude + Interest Mapping",
            ],
        ),
        (
            "Assessment Levels",
            [
                "3 stages: Grades 6–8, 9–10, 11–12",
                "Class 8+, segmented by age",
                "All ages",
                "Youth-focused, flexible",
                "Grades 8+ and graduates",
                "High school to professionals",
                "Grades 9–12 (via TALASH)",
            ],
        ),
        (
            "Report Format",
            [
                "PDF + video explainers + action points for parents & students",
                "Graphical PDF + trait strategies",
                "Text-based PDF",
                "Counsellor-led interpretation",
                "PDF + counsellor sessions",
                "Personality report + career clusters",
                "Career Cards + aptitude scores",
            ],
        ),
        (
            "Career Recommendations",
            [
                "Top 3 clusters with action tasks",
                "Top 5 careers with fitment logic",
                "Career suggestions + application support",
                "Career paths based on strengths",
                "Top 10 matches based on interest",
                "Cluster-based suggestions",
                "Aptitude-interest based suggestions",
            ],
        ),
    ]

    y = TABLE_TOP + HEADER_H + Inches(0.1)
    row_h = Inches(1.15)
    for label, vals in rows:
        draw_table_row(s, y, label, vals, row_h=row_h)
        y += row_h + Inches(0.08)

    # Bottom call-out strip
    add_round(s, Inches(0.6), Inches(6.85), Inches(12.13), Inches(0.45), EMERALD_50, radius=0.4)
    add_text(s, Inches(0.85), Inches(6.85), Inches(11.6), Inches(0.45),
             "Career-9 is the only platform pairing a 6-dimension framework with stage-wise assessments AND parent action plans.",
             size=11, bold=True, color=EMERALD_700,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


# ===================================================================
# Slide 4 — Engagement & Trust
# ===================================================================
def slide_compare_engagement():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, SLATE_50)
    add_page_header(s, "PART 2 OF 2  •  ENGAGEMENT, TECH & TRUST",
                    "What happens after the assessment ends")

    draw_table_header(s, TABLE_TOP)

    rows = [
        (
            "Parental Involvement",
            [
                "Written plan in report",
                "Webinars, sessions, blog resources",
                "Encouraged via sessions",
                "In-session mentoring",
                "Session-based guidance",
                "Personality insights for parents",
                "School-based mentoring & teacher training",
            ],
        ),
        (
            "Post-Assessment Support",
            [
                "90-day LMS habit challenges",
                "Nexus AI, virtual internships, coaching",
                "Career library, mentorship",
                "Strengths-based mentoring",
                "Placement support",
                "Certified counsellor interpretation",
                "Life skills modules, SEL integration",
            ],
        ),
        (
            "Technology Integration",
            [
                "LMS for habit-building",
                "AI-powered personalisation (Nexus)",
                "Web-based dashboard",
                "Phone / web mentoring",
                "Online platform + counsellor access",
                "Multilingual platform",
                "Integrated with NCERT digital tools",
            ],
        ),
        (
            "Scientific Validation",
            [
                "Peer-reviewed (IJIP); endorsed by NIMHANS, UGC, MIT",
                "Peer-reviewed; recognised by Ministry of Science & Tech",
                "Positive reviews; no published benchmarks",
                "Based on positive psychology; no formal validation",
                "Claimed 85% reliability",
                "Statistically validated; Indian context",
                "NCERT-developed; aligned with NEP 2020",
            ],
        ),
    ]

    y = TABLE_TOP + HEADER_H + Inches(0.1)
    row_h = Inches(1.15)
    for label, vals in rows:
        draw_table_row(s, y, label, vals, row_h=row_h)
        y += row_h + Inches(0.08)

    add_round(s, Inches(0.6), Inches(6.85), Inches(12.13), Inches(0.45), EMERALD_50, radius=0.4)
    add_text(s, Inches(0.85), Inches(6.85), Inches(11.6), Inches(0.45),
             "Career-9 turns insight into behaviour: a 90-day LMS habit programme is unique in this list.",
             size=11, bold=True, color=EMERALD_700,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


# ===================================================================
# Slide 5 — The Career-9 Edge
# ===================================================================
def slide_career9_edge():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, SLATE_50)
    add_page_header(s, "THE CAREER-9 EDGE", "Reach, trust and accessibility at a glance")

    # Three highlight cards (Validation / Languages / Cost)
    cards = [
        ("🏛", "Scientific Validation",
         "Peer-reviewed in IJIP and endorsed by NIMHANS, UGC and MIT — independent academic credibility most competitors lack."),
        ("🌐", "Language Accessibility",
         "Available in English, Hindi and Marathi today — built for India's multilingual classrooms."),
        ("🤝", "Cost & Accessibility",
         "Paid model with school partnerships — institution-friendly pricing that scales across grades."),
    ]

    card_w = Inches(3.95)
    card_h = Inches(2.6)
    gap = Inches(0.2)
    total_w = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2
    y = Inches(2.0)

    for i, (icon, title, body) in enumerate(cards):
        x = start_x + (card_w + gap) * i

        # Card
        add_round(s, x, y, card_w, card_h, WHITE, radius=0.06)
        # Left accent stripe
        add_rect(s, x, y, Inches(0.12), card_h, EMERALD_500)

        # Icon circle
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.4), y + Inches(0.35),
                                    Inches(0.85), Inches(0.85))
        circle.fill.solid(); circle.fill.fore_color.rgb = EMERALD_50
        circle.line.color.rgb = EMERALD_500; circle.line.width = Pt(1.5)
        circle.shadow.inherit = False
        add_text(s, x + Inches(0.4), y + Inches(0.35), Inches(0.85), Inches(0.85),
                 icon, size=24, color=EMERALD_700,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Title
        add_text(s, x + Inches(1.4), y + Inches(0.45), card_w - Inches(1.6), Inches(0.55),
                 title, size=16, bold=True, color=SLATE_900)

        # Body
        add_text(s, x + Inches(0.4), y + Inches(1.4), card_w - Inches(0.6), card_h - Inches(1.6),
                 body, size=12, color=SLATE_600, line_spacing=1.4)

    # Closing summary banner
    banner_y = Inches(5.05)
    add_round(s, Inches(0.6), banner_y, Inches(12.13), Inches(1.7), SLATE_900, radius=0.05)
    # Emerald accent strip
    add_rect(s, Inches(0.6), banner_y, Inches(0.12), Inches(1.7), EMERALD_500)

    add_text(s, Inches(1.0), banner_y + Inches(0.25), Inches(11), Inches(0.4),
             "BOTTOM LINE", size=11, bold=True, color=EMERALD_500)
    add_text(s, Inches(1.0), banner_y + Inches(0.55), Inches(11.5), Inches(0.6),
             "Holistic profiling. Action-oriented reports. Real validation.",
             size=22, bold=True, color=WHITE)
    add_text(s, Inches(1.0), banner_y + Inches(1.1), Inches(11.5), Inches(0.5),
             "Career-9 is the only platform combining a 6-dimension framework, family-action plans, and a 90-day behaviour-change LMS.",
             size=12, color=SLATE_300, line_spacing=1.4)


# ===================================================================
# Build all slides
# ===================================================================
def main():
    slide_title()
    slide_why_career9()
    slide_compare_methodology()
    slide_compare_engagement()
    slide_career9_edge()

    total = len(prs.slides)
    # Add footers on slides 2..N (skip the title)
    for idx, slide in enumerate(prs.slides, start=1):
        if idx == 1:
            continue
        add_footer(slide, idx, total)

    out = "Career-9-Comparison.pptx"
    prs.save(out)
    print(f"Wrote {out} ({total} slides)")


if __name__ == "__main__":
    main()
